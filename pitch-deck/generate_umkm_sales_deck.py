import sys
import os
import subprocess

# Ensure python-pptx is installed
try:
    import pptx
except ImportError:
    print("python-pptx not found. Installing...")
    subprocess.run([sys.executable, "-m", "pip", "install", "--user", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# --------------------------------------------------------------------------
# DESIGN SYSTEM VARIABLES (Bestari POS Branding - Slate Teal & Tosca)
# --------------------------------------------------------------------------
COLOR_BG = RGBColor(5, 22, 25)          # Dark Slate Teal (#051619)
COLOR_CARD_BG = RGBColor(8, 33, 37)     # Deep Teal Card (#082125)
COLOR_PRIMARY = RGBColor(13, 148, 136)  # Primary Tosca (#0d9488)
COLOR_LIGHT = RGBColor(45, 212, 191)    # Light Tosca (#2dd4bf)
COLOR_DARK = RGBColor(15, 118, 110)     # Dark Tosca (#0f766e)
COLOR_WHITE = RGBColor(255, 255, 255)   # White text
COLOR_MUTED = RGBColor(148, 163, 184)   # Gray text (#94a3b8)
COLOR_RED = RGBColor(239, 68, 68)       # Red for problems
COLOR_GREEN = RGBColor(16, 185, 129)    # Green for solutions/benefits

FONT_HEADING = "Arial"
FONT_BODY = "Arial"

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --------------------------------------------------------------------------
# UTILITY FUNCTIONS
# --------------------------------------------------------------------------
def apply_slide_bg(slide):
    # Cover background with deep dark slate teal
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG
    bg.line.fill.background()
    return bg

def add_slide_header(slide, tag, title):
    apply_slide_bg(slide)
    
    # Tag Category
    tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.0), Inches(0.4))
    tf = tag_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = tag.upper()
    p.font.name = FONT_HEADING
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLOR_LIGHT
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(10.0), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_HEADING
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Divider Line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_PRIMARY
    line.line.fill.background()

def add_text_box(slide, left, top, width, height, text, size=12, bold=False, italic=False, color=COLOR_WHITE, wrap=True, align=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_BODY
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return txBox

def add_card(slide, left, top, width, height, title, body_paragraphs, accent_color=COLOR_PRIMARY):
    # Rounded Card Background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    card.line.color.rgb = accent_color
    card.line.width = Pt(1.5)
    
    # Card Title
    title_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_HEADING
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Card Body Content
    body_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.7), width - Inches(0.4), height - Inches(0.9))
    tf_body = body_box.text_frame
    tf_body.word_wrap = True
    tf_body.margin_left = tf_body.margin_right = tf_body.margin_top = tf_body.margin_bottom = Inches(0)
    
    for idx, text in enumerate(body_paragraphs):
        p = tf_body.add_paragraph() if idx > 0 else tf_body.paragraphs[0]
        p.text = text
        p.font.name = FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_MUTED
        p.space_after = Pt(6)

def add_table_custom(slide, left, top, width, height, headers, rows):
    rows_cnt = len(rows) + 1
    cols_cnt = len(headers)
    table_shape = slide.shapes.add_table(rows_cnt, cols_cnt, left, top, width, height)
    table = table_shape.table
    
    # Stylize headers
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.font.name = FONT_HEADING
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.alignment = PP_ALIGN.CENTER
        
    # Stylize rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            cell.fill.solid()
            # Alternate row backgrounds
            if row_idx % 2 == 0:
                cell.fill.fore_color.rgb = COLOR_CARD_BG
            else:
                cell.fill.fore_color.rgb = RGBColor(10, 42, 47)
            p = cell.text_frame.paragraphs[0]
            p.font.name = FONT_BODY
            p.font.size = Pt(11)
            p.font.color.rgb = COLOR_WHITE
            if col_idx == 0:
                p.alignment = PP_ALIGN.LEFT
            else:
                p.alignment = PP_ALIGN.CENTER

# --------------------------------------------------------------------------
# SLIDE 1: COVER SLIDE
# --------------------------------------------------------------------------
slide_1 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_bg(slide_1)

# Large decorative glow card behind text
glow_rect = slide_1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(5.8), Inches(7.5))
glow_rect.fill.solid()
glow_rect.fill.fore_color.rgb = COLOR_CARD_BG
glow_rect.line.fill.background()

# Cover Title & Subtitle
add_text_box(slide_1, Inches(0.8), Inches(1.8), Inches(8.0), Inches(0.5), "APLIKASI KASIR WIRAUSAHA", size=11, bold=True, color=COLOR_LIGHT)
add_text_box(slide_1, Inches(0.8), Inches(2.2), Inches(8.0), Inches(1.2), "Bestari POS", size=54, bold=True, color=COLOR_WHITE)
add_text_box(slide_1, Inches(0.8), Inches(3.4), Inches(8.0), Inches(1.0), 
             "Solusi Kasir Pintar Offline-First. Jualan Lancar Tanpa Takut Internet Macet, Kelola Stok & Laporan Langsung di Genggaman.", size=15, bold=False, color=COLOR_MUTED)

# Tagline Box
add_text_box(slide_1, Inches(0.8), Inches(4.5), Inches(8.0), Inches(0.6), 
             '"Transaksi Jalan Terus, Staf Nyaman, Omzet Terpantau Aman!"', size=13, italic=True, color=COLOR_LIGHT)

# Footer info on Cover
meta_y = Inches(5.6)
add_text_box(slide_1, Inches(0.8), meta_y, Inches(2.5), Inches(0.3), "UNTUK INFORMASI & DEMO", size=9, bold=True, color=COLOR_MUTED)
add_text_box(slide_1, Inches(0.8), meta_y + Inches(0.25), Inches(2.5), Inches(0.4), "info@bestaripos.com", size=13, bold=True, color=COLOR_WHITE)

add_text_box(slide_1, Inches(3.8), meta_y, Inches(3.0), Inches(0.3), "PROGRAM KEMITRAAN UMKM", size=9, bold=True, color=COLOR_MUTED)
add_text_box(slide_1, Inches(3.8), meta_y + Inches(0.25), Inches(3.0), Inches(0.4), "Gratis Uji Coba 14 Hari", size=13, bold=True, color=COLOR_LIGHT)

# Cover Illustration shape (mockup box showing a friendly menu checkout)
ill_card = slide_1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.5), Inches(4.0), Inches(4.5))
ill_card.fill.solid()
ill_card.fill.fore_color.rgb = COLOR_CARD_BG
ill_card.line.color.rgb = COLOR_PRIMARY
ill_card.line.width = Pt(2)

add_text_box(slide_1, Inches(8.8), Inches(1.8), Inches(3.4), Inches(0.4), "Bestari Kasir App", size=14, bold=True, color=COLOR_WHITE)
# Draw visual horizontal line inside mockup
line_inner = slide_1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.8), Inches(2.2), Inches(3.4), Inches(0.015))
line_inner.fill.solid()
line_inner.fill.fore_color.rgb = COLOR_DARK
line_inner.line.fill.background()

# Inner Mockup List Items
mock_items = [
    ("☕ Kopi Susu Aren", "1x Rp 18.000"),
    ("🥐 Croissant Butter", "2x Rp 44.000"),
    ("🍵 Matcha Latte", "1x Rp 24.000")
]
for idx, (name, val) in enumerate(mock_items):
    item_y = Inches(2.4) + Inches(idx * 0.7)
    add_text_box(slide_1, Inches(8.8), item_y, Inches(2.0), Inches(0.3), name, size=11, bold=True)
    add_text_box(slide_1, Inches(8.8), item_y + Inches(0.25), Inches(2.0), Inches(0.3), val, size=10, color=COLOR_MUTED)
    # Circle indicator for checkout status
    circle = slide_1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.8), item_y + Inches(0.1), Inches(0.25), Inches(0.25))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLOR_PRIMARY
    circle.line.fill.background()

# Total Checkout Area
chk_y = Inches(4.7)
add_text_box(slide_1, Inches(8.8), chk_y, Inches(2.0), Inches(0.4), "Total: Rp 86.000", size=14, bold=True, color=COLOR_LIGHT)
pay_btn = slide_1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), chk_y + Inches(0.4), Inches(3.4), Inches(0.5))
pay_btn.fill.solid()
pay_btn.fill.fore_color.rgb = COLOR_PRIMARY
pay_btn.line.fill.background()
# Text inside button
add_text_box(slide_1, Inches(8.8), chk_y + Inches(0.45), Inches(3.4), Inches(0.4), "BAYAR (BISA OFFLINE)", size=10, bold=True, align=PP_ALIGN.CENTER)

# --------------------------------------------------------------------------
# SLIDE 2: THE PAIN POINTS OF UMKM
# --------------------------------------------------------------------------
slide_2 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide_2, "Masalah Wirausaha", "Pernahkah Bisnis Anda Mengalami Masalah Ini?")

# Grid layout with 3 problem cards
add_card(slide_2, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.5), 
         "1. Internet Mati = Kasir Macet", 
         [
             "Saat koneksi Wi-Fi bermasalah atau sinyal provider melemah, kasir berbasis cloud biasa langsung hang.",
             "Staf panik, terjadi antrean panjang di meja kasir, dan pembeli yang kesal akhirnya pergi ke kompetitor.",
             "Operasional harian lumpuh total!"
         ], accent_color=COLOR_RED)

add_card(slide_2, Inches(4.8), Inches(2.0), Inches(3.6), Inches(4.5), 
         "2. Stok Bahan Baku Sering Bocor", 
         [
             "Bahan mahal seperti susu, biji kopi, daging, atau cup saji sering habis tidak wajar tanpa alasan yang jelas.",
             "Sangat sulit menghitung secara akurat kecocokan stok fisik dengan porsi menu terjual.",
             "Bocornya stok memotong profit bulanan Anda secara diam-diam."
         ], accent_color=COLOR_RED)

add_card(slide_2, Inches(8.8), Inches(2.0), Inches(3.6), Inches(4.5), 
         "3. Biaya Bulanan & Alat Mahal", 
         [
             "Banyak vendor kasir mengenakan biaya bulanan yang mahal per perangkat yang terhubung.",
             "Ada juga biaya tambahan untuk fitur wajib seperti modul dapur atau manajemen stok.",
             "Beberapa bahkan mengharuskan Anda membeli alat mesin kasir khusus mereka."
         ], accent_color=COLOR_RED)

# --------------------------------------------------------------------------
# SLIDE 3: THE SOLUTION - KASIR OFFLINE-FIRST
# --------------------------------------------------------------------------
slide_3 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide_3, "Solusi Kami", "Bestari POS: Kasir Lancar Bebas Gangguan Internet")

# Left Column (Content bullet list)
left_x = Inches(0.8)
add_text_box(slide_3, left_x, Inches(2.0), Inches(5.8), Inches(1.0), 
             "Aplikasi Kasir Bestari didesain dengan teknologi modern offline-first yang bekerja mandiri di perangkat kasir Anda. Internet kencang? Bagus. Internet mati? Tidak masalah!", 
             size=14, color=COLOR_WHITE)

bullet_y = Inches(3.2)
# Bullet 1
circle1 = slide_3.shapes.add_shape(MSO_SHAPE.OVAL, left_x, bullet_y, Inches(0.3), Inches(0.3))
circle1.fill.solid()
circle1.fill.fore_color.rgb = COLOR_GREEN
circle1.line.fill.background()
add_text_box(slide_3, left_x + Inches(0.5), bullet_y - Inches(0.05), Inches(5.3), Inches(1.0), 
             "Toleran Terhadap Gangguan Internet\nSemua transaksi kasir, input menu, diskon, dan cetak struk berjalan 100% normal tanpa jaringan internet.", size=12, color=COLOR_MUTED)

# Bullet 2
bullet_y2 = bullet_y + Inches(1.1)
circle2 = slide_3.shapes.add_shape(MSO_SHAPE.OVAL, left_x, bullet_y2, Inches(0.3), Inches(0.3))
circle2.fill.solid()
circle2.fill.fore_color.rgb = COLOR_GREEN
circle2.line.fill.background()
add_text_box(slide_3, left_x + Inches(0.5), bullet_y2 - Inches(0.05), Inches(5.3), Inches(1.0), 
             "Auto Sync ke Cloud Secara Aman\nBegitu internet terhubung kembali, aplikasi secara cerdas mensinkronkan transaksi ke server awan (cloud) di latar belakang.", size=12, color=COLOR_MUTED)

# Bullet 3
bullet_y3 = bullet_y2 + Inches(1.1)
circle3 = slide_3.shapes.add_shape(MSO_SHAPE.OVAL, left_x, bullet_y3, Inches(0.3), Inches(0.3))
circle3.fill.solid()
circle3.fill.fore_color.rgb = COLOR_GREEN
circle3.line.fill.background()
add_text_box(slide_3, left_x + Inches(0.5), bullet_y3 - Inches(0.05), Inches(5.3), Inches(1.0), 
             "Hemat Data Seluler & Hardware\nTidak memerlukan data internet berkapasitas besar. Dapat berjalan di HP kasir yang Anda miliki sekarang tanpa perlu beli alat baru.", size=12, color=COLOR_MUTED)

# Right Column: Visual Diagram Cards
right_x = Inches(7.2)
add_card(slide_3, right_x, Inches(2.2), Inches(5.3), Inches(1.2), 
         "💻 Tablet / HP Kasir Anda", 
         ["Menyimpan data penjualan, katalog menu, diskon secara aman langsung di memori browser perangkat (Offline Mode)."])

# Connector Arrow Shape
arrow_shape = slide_3.shapes.add_shape(MSO_SHAPE.UP_DOWN_ARROW, right_x + Inches(2.2), Inches(3.6), Inches(0.8), Inches(0.6))
arrow_shape.fill.solid()
arrow_shape.fill.fore_color.rgb = COLOR_LIGHT
arrow_shape.line.fill.background()
add_text_box(slide_3, right_x + Inches(3.2), Inches(3.7), Inches(2.0), Inches(0.4), "Otomatis Sinkron", size=10, bold=True, color=COLOR_LIGHT)

add_card(slide_3, right_x, Inches(4.4), Inches(5.3), Inches(1.2), 
         "☁️ Database Cloud Bestari", 
         ["Mengamankan backup data transaksi harian, menyajikan dashboard grafis penjualan, dan memantau multi-outlet wirausaha."])

# --------------------------------------------------------------------------
# SLIDE 4: KASIR INSTAN & ANTREAN DAPUR
# --------------------------------------------------------------------------
slide_4 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide_4, "Fitur Unggulan 1", "Transaksi Kasir Instan & Antrean Dapur Teratur")

# Left side description
add_text_box(slide_4, Inches(0.8), Inches(2.0), Inches(5.5), Inches(1.0),
             "Proses transaksi kasir secepat kilat dengan antrean yang terpantau rapi dari depan (kasir) ke belakang (barista/koki).", size=14, color=COLOR_WHITE)

add_card(slide_4, Inches(0.8), Inches(3.0), Inches(5.5), Inches(1.8),
         "🛒 Kasir Super Cepat",
         [
             "• Menu instan dengan variasi rasa (add-on) dan modifikasi porsi.",
             "• Input diskon potongan harga langsung atau persentase.",
             "• Pencatatan jenis pembayaran tunai maupun e-wallet/cashless.",
             "• Cetak struk struk fisik atau kirim struk digital lewat WhatsApp."
         ])

add_card(slide_4, Inches(0.8), Inches(5.0), Inches(5.5), Inches(1.8),
         "👥 Layar Antrean Dapur (Kitchen Display)",
         [
             "• Hubungkan tablet/HP di dapur untuk menerima pesanan real-time.",
             "• Singkirkan kertas struk cetak yang gampang hilang atau basah.",
             "• Status visual pesanan yang jelas: Menunggu -> Diproses -> Selesai.",
             "• Mengurangi kesalahan staf dapur dalam menyajikan menu."
         ])

# Right side: Antrean visual design mock card
kitchen_card = slide_4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.0), Inches(2.2), Inches(5.5), Inches(4.5))
kitchen_card.fill.solid()
kitchen_card.fill.fore_color.rgb = COLOR_CARD_BG
kitchen_card.line.color.rgb = COLOR_LIGHT
kitchen_card.line.width = Pt(1.5)

add_text_box(slide_4, Inches(7.3), Inches(2.5), Inches(4.9), Inches(0.4), "Visual Antrean Dapur (Dapur View)", size=14, bold=True, color=COLOR_WHITE)

# 3 horizontal layout task columns representing status queue
q_items = [
    ("MENUNGGU", "Kopi Susu Aren - 1x\nMatcha Latte - 1x", COLOR_RED),
    ("DIPROSES", "Choco Lava - 2x\nEspresso - 1x", COLOR_PRIMARY),
    ("SIAP SAJI", "Croissant Butter - 2x\nIce Lemon Tea - 1x", COLOR_GREEN)
]

for idx, (title, content, color) in enumerate(q_items):
    col_x = Inches(7.3) + Inches(idx * 1.65)
    
    # Status Header
    hdr = slide_4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col_x, Inches(3.1), Inches(1.55), Inches(0.4))
    hdr.fill.solid()
    hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    add_text_box(slide_4, col_x, Inches(3.15), Inches(1.55), Inches(0.3), title, size=9, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    
    # Status body card
    body = slide_4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, col_x, Inches(3.6), Inches(1.55), Inches(2.6))
    body.fill.solid()
    body.fill.fore_color.rgb = COLOR_BG
    body.line.color.rgb = COLOR_DARK
    body.line.width = Pt(1)
    
    add_text_box(slide_4, col_x + Inches(0.08), Inches(3.7), Inches(1.4), Inches(2.4), content, size=9, color=COLOR_MUTED)

# --------------------------------------------------------------------------
# SLIDE 5: INVENTORY / INGREDIENT DEDUCTION
# --------------------------------------------------------------------------
slide_5 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide_5, "Fitur Unggulan 2", "Kontrol Bahan Baku Otomatis & Stop Bocor!")

# Left design columns
add_card(slide_5, Inches(0.8), Inches(2.2), Inches(5.5), Inches(4.5),
         "Mengapa Manajemen Bahan Baku Bestari POS Penting?",
         [
             "Aplikasi kami memiliki modul kalkulator resep terintegrasi.",
             "Anda dapat menentukan formula resep untuk setiap menu. Misalnya, 1 gelas Kopi Susu Aren membutuhkan:\n- Biji Kopi Espresso: 15 gram\n- Susu Fresh Milk: 100 ml\n- Gula Aren Cair: 20 ml\n- Cup & Sedotan: 1 unit",
             "Setiap kali kasir mengetuk tombol bayar untuk menu tersebut, stok bahan baku Anda langsung terpotong secara otomatis di gudang penyimpanan.",
             "• Ketahui stok menipis secara instan.",
             "• Hindari kecurangan atau hilangnya bahan tanpa pencatatan.",
             "• Belanja bahan baku terjadwal sesuai laporan sisa stok."
         ])

# Right design: visual representation of ingredient deduction
right_x = Inches(7.0)
add_text_box(slide_5, right_x, Inches(2.0), Inches(5.5), Inches(0.4), "Simulasi Pengurangan Bahan Baku Otomatis", size=13, bold=True, color=COLOR_LIGHT)

simulations = [
    ("☕ Transaksi Kasir", "Pembeli memesan 2x Kopi Susu Aren", COLOR_PRIMARY),
    ("🧮 Hitung Resep", "Sistem mengalikan porsi (2x) dengan resep terdaftar", COLOR_PRIMARY),
    ("📉 Stok Terpotong", "Biji Kopi -30g, Fresh Milk -200ml, Cup -2pcs", COLOR_GREEN),
    ("⚠️ Peringatan Stok", "Stok Fresh Milk menipis! (Sisa 1.2 Liter)", COLOR_RED)
]

for idx, (title, desc, color) in enumerate(simulations):
    sim_y = Inches(2.6) + Inches(idx * 1.1)
    
    # Step indicator
    step_circle = slide_5.shapes.add_shape(MSO_SHAPE.OVAL, right_x, sim_y + Inches(0.1), Inches(0.6), Inches(0.6))
    step_circle.fill.solid()
    step_circle.fill.fore_color.rgb = COLOR_CARD_BG
    step_circle.line.color.rgb = color
    step_circle.line.width = Pt(2)
    add_text_box(slide_5, right_x, sim_y + Inches(0.2), Inches(0.6), Inches(0.4), f"0{idx+1}", size=11, bold=True, color=color, align=PP_ALIGN.CENTER)
    
    # Text explanation beside step
    add_text_box(slide_5, right_x + Inches(0.8), sim_y, Inches(4.7), Inches(0.3), title, size=12, bold=True, color=COLOR_WHITE)
    add_text_box(slide_5, right_x + Inches(0.8), sim_y + Inches(0.25), Inches(4.7), Inches(0.6), desc, size=10, color=COLOR_MUTED)
    
    # Small line connector
    if idx < 3:
        conn = slide_5.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x + Inches(0.28), sim_y + Inches(0.7), Inches(0.04), Inches(0.4))
        conn.fill.solid()
        conn.fill.fore_color.rgb = COLOR_DARK
        conn.line.fill.background()

# --------------------------------------------------------------------------
# SLIDE 6: REPORTS & OWNER MOBILE ACCESS
# --------------------------------------------------------------------------
slide_6 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide_6, "Fitur Unggulan 3", "Laporan Keuangan Real-Time & Akses Owner Aman")

# 4 feature cards
add_card(slide_6, Inches(0.8), Inches(2.2), Inches(2.7), Inches(4.3), 
         "📊 Laporan Omzet", 
         [
             "Lihat total penjualan harian, mingguan, atau bulanan secara grafis.",
             "Ketahui margin profit bersih kotor secara otomatis."
         ])

add_card(slide_6, Inches(3.8), Inches(2.2), Inches(2.7), Inches(4.3), 
         "📈 Produk Terlaris", 
         [
             "Sistem menampilkan produk yang paling laku terjual.",
             "Hapus menu yang tidak laku dan optimalkan stok menu terlaris."
         ])

add_card(slide_6, Inches(6.8), Inches(2.2), Inches(2.7), Inches(4.3), 
         "🔑 Keamanan PIN", 
         [
             "Batasi akses staf kasir Anda.",
             "Kasir hanya bisa memasukkan order, sedangkan data laporan keuangan dikunci hanya untuk owner."
         ])

add_card(slide_6, Inches(9.8), Inches(2.2), Inches(2.7), Inches(4.3), 
         "☁️ Pantau dari Jauh", 
         [
             "Laporan terkirim ke cloud secara berkala.",
             "Owner dapat melihat performa kedai dari HP di rumah tanpa mengganggu kasir."
         ])

# --------------------------------------------------------------------------
# SLIDE 7: DEVICE COMPATIBILITY & PWA
# --------------------------------------------------------------------------
slide_7 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide_7, "Kemudahan Instalasi", "Cukup Gunakan HP, Tablet, atau Laptop yang Anda Punya")

# Left Column (Device images representations)
left_x = Inches(0.8)
add_text_box(slide_7, left_x, Inches(2.0), Inches(5.5), Inches(0.8), 
             "Tidak ada keharusan membeli perangkat kasir khusus yang mahal. Aplikasi Bestari POS siap dipasang di alat apa pun milik Anda saat ini.", 
             size=14, color=COLOR_WHITE)

add_card(slide_7, left_x, Inches(3.0), Inches(5.5), Inches(3.5), 
         "💻 Kebebasan Memilih Perangkat", 
         [
             "• Tablet (iPad / Android Tablet): Sangat ideal untuk pajangan kasir meja depan yang rapi.",
             "• Handphone (Android / iOS): Kasir mobile, koki di dapur, atau pelayan di meja bisa memesan langsung dari HP.",
             "• Laptop & PC: Sangat nyaman untuk outlet dengan volume transaksi besar atau printer struk kabel.",
             "• Printer Struk Bluetooth & USB: Kompatibel penuh dengan berbagai merek printer struk termal di pasaran."
         ])

# Right Column (PWA Installation steps)
right_x = Inches(7.0)
add_text_box(slide_7, right_x, Inches(2.0), Inches(5.5), Inches(0.4), "Langkah Instalasi Instan (PWA)", size=13, bold=True, color=COLOR_LIGHT)

steps = [
    ("1. Buka Link Aplikasi", "Buka browser Chrome/Safari di HP/Tablet, lalu ketik alamat website POS."),
    ("2. Klik 'Install App'", "Klik tombol instalasi aplikasi yang muncul langsung di layar browser Anda."),
    ("3. Siap Dipakai!", "Ikon aplikasi kasir akan muncul di Home Screen. Buka, masukkan PIN, dan langsung jualan!")
]

for idx, (title, desc) in enumerate(steps):
    step_y = Inches(2.6) + Inches(idx * 1.4)
    
    # Card representing PWA Step
    step_card = slide_7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, right_x, step_y, Inches(5.5), Inches(1.15))
    step_card.fill.solid()
    step_card.fill.fore_color.rgb = COLOR_CARD_BG
    step_card.line.color.rgb = COLOR_PRIMARY
    step_card.line.width = Pt(1)
    
    add_text_box(slide_7, right_x + Inches(0.2), step_y + Inches(0.15), Inches(5.1), Inches(0.3), title, size=12, bold=True, color=COLOR_LIGHT)
    add_text_box(slide_7, right_x + Inches(0.2), step_y + Inches(0.45), Inches(5.1), Inches(0.6), desc, size=10, color=COLOR_MUTED)

# --------------------------------------------------------------------------
# SLIDE 8: PRICING PLANS
# --------------------------------------------------------------------------
slide_8 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide_8, "Paket Harga", "Pilihan Paket Langganan Hemat & Tanpa Biaya Tersembunyi")

# Grid of 3 pricing plans
add_card(slide_8, Inches(0.8), Inches(2.2), Inches(3.6), Inches(4.4), 
         "Starter (Lokal Mandiri)", 
         [
             "Rp 149.000 / Bulan",
             "------------------------",
             "Cocok untuk kedai kecil mandiri.",
             "",
             "✔ 1 Perangkat kasir aktif.",
             "✔ Fitur kasir lengkap.",
             "✔ Pilihan Add-on / varian.",
             "✔ Penyimpanan lokal di browser.",
             "✔ Ekspor laporan harian ke Excel.",
             "✔ Berjalan 100% offline."
         ])

# Middle card represents the recommended plan (using secondary light accent)
add_card(slide_8, Inches(4.8), Inches(2.2), Inches(3.8), Inches(4.6), 
         "Professional (Rekomendasi)", 
         [
             "Rp 299.000 / Bulan",
             "------------------------",
             "Untuk bisnis berkembang & multi-outlet.",
             "",
             "✔ Hubungkan banyak kasir & dapur.",
             "✔ Cloud backup otomatis ke Supabase.",
             "✔ Modul resep & stok bahan gudang.",
             "✔ Dashboard owner (laporan real-time).",
             "✔ Pembaruan fitur otomatis berkala.",
             "✔ Dukungan prioritas tim support."
         ], accent_color=COLOR_LIGHT)

add_card(slide_8, Inches(8.9), Inches(2.2), Inches(3.6), Inches(4.4), 
         "Custom / Franchise", 
         [
             "Hubungi Tim Sales",
             "------------------------",
             "Untuk jaringan usaha waralaba / franchise.",
             "",
             "✔ Nama domain & logo aplikasi kustom.",
             "✔ Server cloud dedicated terpisah.",
             "✔ Custom integrasi pembayaran QRIS.",
             "✔ Hak milik data penjualan penuh.",
             "✔ Pendampingan setup di awal.",
             "✔ SLA Bantuan Teknis 24/7."
         ])

# --------------------------------------------------------------------------
# SLIDE 9: COMPARISON VS COMPETITORS
# --------------------------------------------------------------------------
slide_9 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide_9, "Mengapa Bestari POS?", "Bandingkan Keunggulan Bestari POS dari Kasir Biasa")

# Text intro
add_text_box(slide_9, Inches(0.8), Inches(1.8), Inches(11.733), Inches(0.4),
             "Berikut perbandingan fitur utama Bestari POS dengan aplikasi kasir cloud pada umumnya:", size=12, color=COLOR_WHITE)

headers = ["Fitur Kasir", "Bestari POS (Offline-First)", "Aplikasi Kasir Cloud Umum", "Mesin Kasir Tradisional"]
rows = [
    ["Konektivitas Internet", "Bekerja normal 100% tanpa internet (Downtime-safe)", "Kasir tidak bisa transaksi saat internet putus", "Offline penuh tapi tidak ada sinkronisasi cloud"],
    ["Biaya Bulanan", "Sangat hemat & transparan (Mulai Rp 149rb/outlet)", "Mahal (sering Rp 300rb-500rb/alat/bulan)", "Sekali bayar tapi beli hardware mahal sekali di awal"],
    ["Kontrol Stok Bahan", "Otomatis memotong bahan resep terperinci (real-time)", "Banyak yang harus bayar modul stok terpisah", "Manual mencocokkan di kertas atau Excel"],
    ["Fleksibilitas Alat", "Kompatibel di HP kasir lama, tablet, maupun laptop", "Wajib spesifikasi tablet baru atau terminal hardware vendor", "Hanya bisa di monitor mesin kasir bawaan pabrik"],
    ["Laporan Owner", "Akses dashboard real-time dari mana saja via cloud", "Bisa dipantau tapi sering terhambat data kasir macet", "Harus datang langsung ke toko untuk melihat laporan"]
]

# Generate custom table
add_table_custom(slide_9, Inches(0.8), Inches(2.3), Inches(11.733), Inches(4.2), headers, rows)

# --------------------------------------------------------------------------
# SLIDE 10: PENUTUP & KONTAK
# --------------------------------------------------------------------------
slide_10 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_bg(slide_10)

# Large accent container box
container = slide_10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
container.fill.solid()
container.fill.fore_color.rgb = COLOR_CARD_BG
container.line.color.rgb = COLOR_PRIMARY
container.line.width = Pt(2)

# Closing Text
add_text_box(slide_10, Inches(2.0), Inches(2.2), Inches(9.333), Inches(0.5), 
             "Siap Mengembangkan Usaha Kuliner Anda?", size=24, bold=True, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)

add_text_box(slide_10, Inches(2.0), Inches(3.0), Inches(9.333), Inches(0.8), 
             "Jangan biarkan masalah internet lambat mengganggu transaksi Anda. Mulai uji coba Bestari POS sekarang juga secara gratis selama 14 hari penuh!", 
             size=14, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

# Contact info dividers
divider = slide_10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.0), Inches(3.9), Inches(5.333), Inches(0.015))
divider.fill.solid()
divider.fill.fore_color.rgb = COLOR_DARK
divider.line.fill.background()

add_text_box(slide_10, Inches(2.0), Inches(4.2), Inches(9.333), Inches(1.2), 
             "Hubungi Tim Sales Bestari POS untuk Jadwal Demo Gratis:\n📞 WhatsApp: +62 812-3456-7890   |   📧 Email: sales@bestaripos.com\n🌐 Website: bestaripos.com   |   📍 Kantor Layanan: Jakarta, Indonesia", 
             size=12, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

# Save the generated presentation
script_dir = os.path.dirname(os.path.abspath(__file__))
project_root = os.path.dirname(script_dir)
output_path = os.path.join(project_root, "Bestari_POS_Sales_Deck_UMKM.pptx")

try:
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")
except Exception as e:
    # Fallback to local script folder if save fails
    output_path = os.path.join(script_dir, "Bestari_POS_Sales_Deck_UMKM.pptx")
    prs.save(output_path)
    print(f"Presentation saved successfully to: {output_path}")
