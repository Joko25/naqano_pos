import sys
import os
import subprocess

# Ensure python-pptx is installed
try:
    import pptx
except ImportError:
    print("python-pptx not found. Installing...")
    subprocess.run([sys.executable, "-m", "pip", "install", "--user", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# --------------------------------------------------------------------------
# DESIGN SYSTEM VARIABLES
# --------------------------------------------------------------------------
COLOR_BG = RGBColor(5, 22, 25)          # Dark Slate Teal (#051619)
COLOR_CARD_BG = RGBColor(8, 33, 37)     # Deep Teal Card (#082125)
COLOR_PRIMARY = RGBColor(13, 148, 136)  # Primary Tosca (#0d9488)
COLOR_LIGHT = RGBColor(45, 212, 191)    # Light Tosca (#2dd4bf)
COLOR_DARK = RGBColor(15, 118, 110)     # Dark Tosca (#0f766e)
COLOR_WHITE = RGBColor(255, 255, 255)   # White text
COLOR_MUTED = RGBColor(148, 163, 184)   # Gray text (#94a3b8)
COLOR_RED = RGBColor(239, 68, 68)       # Red for error/problems
COLOR_GREEN = RGBColor(16, 185, 129)    # Green for success/solutions

FONT_HEADING = "Arial"
FONT_BODY = "Arial"

# Initialize Presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# --------------------------------------------------------------------------
# UTILITY FUNCTIONS
# --------------------------------------------------------------------------
def apply_slide_bg(slide):
    # Cover background with deep dark slate teal
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLOR_BG
    bg.line.fill.background()
    return bg

def add_slide_header(slide, tag, title):
    apply_slide_bg(slide)
    
    # Tag Category
    tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(8.0), Inches(0.4))
    tf = tag_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = tag.upper()
    p.font.name = FONT_HEADING
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = COLOR_LIGHT
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(10.0), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_HEADING
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Divider Line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_PRIMARY
    line.line.fill.background()

def add_text_box(slide, left, top, width, height, text, size=12, bold=False, italic=False, color=COLOR_WHITE, wrap=True, align=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_BODY
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return txBox

def add_card(slide, left, top, width, height, title, body_paragraphs, accent_color=COLOR_PRIMARY):
    # Rounded Card Background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_CARD_BG
    card.line.color.rgb = accent_color
    card.line.width = Pt(1.5)
    
    # Card Title
    title_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Inches(0)
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_HEADING
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    
    # Card Body Content
    body_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.7), width - Inches(0.4), height - Inches(0.9))
    tf_body = body_box.text_frame
    tf_body.word_wrap = True
    tf_body.margin_left = tf_body.margin_right = tf_body.margin_top = tf_body.margin_bottom = Inches(0)
    
    for idx, text in enumerate(body_paragraphs):
        p = tf_body.add_paragraph() if idx > 0 else tf_body.paragraphs[0]
        p.text = text
        p.font.name = FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_MUTED
        p.space_after = Pt(6)

# --------------------------------------------------------------------------
# SLIDE 1: COVER SLIDE
# --------------------------------------------------------------------------
slide_1 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_bg(slide_1)

# Large decorative glow card behind text
glow_rect = slide_1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(5.5), Inches(7.5))
glow_rect.fill.solid()
glow_rect.fill.fore_color.rgb = COLOR_CARD_BG
glow_rect.line.fill.background()

# Cover Title
add_text_box(slide_1, Inches(0.8), Inches(1.8), Inches(8.0), Inches(0.5), "INVESTOR PITCH DECK", size=11, bold=True, color=COLOR_LIGHT)
add_text_box(slide_1, Inches(0.8), Inches(2.2), Inches(8.0), Inches(1.2), "Bestari POS", size=54, bold=True, color=COLOR_WHITE)
add_text_box(slide_1, Inches(0.8), Inches(3.4), Inches(8.0), Inches(1.0), 
             "Revolusi Point of Sale Offline-First untuk Digitalisasi Operasional F&B dan UMKM Kuliner.", size=16, bold=False, color=COLOR_MUTED)

# Tagline Box
add_text_box(slide_1, Inches(0.8), Inches(4.5), Inches(8.0), Inches(0.6), 
             '"Transaksi Jalan Terus, Tanpa Khawatir Internet Putus."', size=13, italic=True, color=COLOR_LIGHT)

# Meta info lines
meta_y = Inches(5.5)
add_text_box(slide_1, Inches(0.8), meta_y, Inches(2.5), Inches(0.3), "PRESENTED BY", size=9, bold=True, color=COLOR_MUTED)
add_text_box(slide_1, Inches(0.8), meta_y + Inches(0.25), Inches(2.5), Inches(0.4), "Tim Pengembang Bestari", size=13, bold=True, color=COLOR_WHITE)

add_text_box(slide_1, Inches(3.8), meta_y, Inches(3.0), Inches(0.3), "TARGET PENDANAAN", size=9, bold=True, color=COLOR_MUTED)
add_text_box(slide_1, Inches(3.8), meta_y + Inches(0.25), Inches(3.0), Inches(0.4), "Rp 2 Miliar (Seed Round)", size=13, bold=True, color=COLOR_LIGHT)

# Cover Illustration shape (mockup box)
ill_card = slide_1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.5), Inches(1.5), Inches(4.0), Inches(4.5))
ill_card.fill.solid()
ill_card.fill.fore_color.rgb = COLOR_CARD_BG
ill_card.line.color.rgb = COLOR_PRIMARY
ill_card.line.width = Pt(2)

add_text_box(slide_1, Inches(8.8), Inches(1.8), Inches(3.4), Inches(0.4), "Bestari Kasir App", size=14, bold=True, color=COLOR_WHITE)
# Draw visual horizontal line inside mockup
line_inner = slide_1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.8), Inches(2.2), Inches(3.4), Inches(0.015))
line_inner.fill.solid()
line_inner.fill.fore_color.rgb = COLOR_DARK
line_inner.line.fill.background()

# Inner Mockup List Items
mock_items = [
    ("☕ Kopi Susu Aren", "1x Rp 18.000"),
    ("🥐 Croissant Butter", "2x Rp 44.000"),
    ("🍵 Matcha Latte", "1x Rp 24.000")
]
for idx, (name, val) in enumerate(mock_items):
    item_y = Inches(2.4) + Inches(idx * 0.7)
    add_text_box(slide_1, Inches(8.8), item_y, Inches(2.0), Inches(0.3), name, size=11, bold=True)
    add_text_box(slide_1, Inches(8.8), item_y + Inches(0.25), Inches(2.0), Inches(0.3), val, size=10, color=COLOR_MUTED)
    # Circle indicator
    circle = slide_1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(11.8), item_y + Inches(0.1), Inches(0.25), Inches(0.25))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLOR_PRIMARY
    circle.line.fill.background()

# Total Checkout Area
chk_y = Inches(4.7)
add_text_box(slide_1, Inches(8.8), chk_y, Inches(2.0), Inches(0.4), "Total: Rp 86.000", size=14, bold=True, color=COLOR_LIGHT)
pay_btn = slide_1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), chk_y + Inches(0.4), Inches(3.4), Inches(0.5))
pay_btn.fill.solid()
pay_btn.fill.fore_color.rgb = COLOR_PRIMARY
pay_btn.line.fill.background()
# Text inside button
add_text_box(slide_1, Inches(8.8), chk_y + Inches(0.45), Inches(3.4), Inches(0.4), "PROSES BAYAR (OFFLINE READY)", size=10, bold=True, align=PP_ALIGN.CENTER)

# --------------------------------------------------------------------------
# SLIDE 2: THE PROBLEM
# --------------------------------------------------------------------------
slide_2 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide_2, "The Problem", "Masalah Utama POS Tradisional di Lapangan")

# Grid layout with 3 problem cards
add_card(slide_2, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.5), 
         "1. Ketergantungan Internet", 
         [
             "Sebagian besar POS modern bergantung penuh pada jaringan cloud.",
             "Saat internet putus atau lambat di jam sibuk, operasional kasir lumpuh total.",
             "Hal ini memicu antrean panjang, kepanikan staf, dan kekecewaan pelanggan."
         ], accent_color=COLOR_RED)

add_card(slide_2, Inches(4.8), Inches(2.0), Inches(3.6), Inches(4.5), 
         "2. Kehilangan Data Penjualan", 
         [
             "Sistem pencatatan manual atau POS offline tradisional rentan mengalami kegagalan hardware lokal.",
             "Sinkronisasi data manual pasca-koneksi mati sering memicu konflik data (conflict) dan tumpang tindih transaksi."
         ], accent_color=COLOR_RED)

add_card(slide_2, Inches(8.8), Inches(2.0), Inches(3.6), Inches(4.5), 
         "3. Beban Biaya Bulanan Tinggi", 
         [
             "Vendor POS cloud membebankan tarif langganan bulanan per perangkat (device) yang memberatkan UMKM.",
             "Seringkali ada biaya tambahan (add-on) terpisah untuk modul stok bahan baku dan antrean dapur."
         ], accent_color=COLOR_RED)

# --------------------------------------------------------------------------
# SLIDE 3: THE SOLUTION
# --------------------------------------------------------------------------
slide_3 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide_3, "The Solution", "Bestari POS: Arsitektur Kasir Offline-First Tangguh")

# Left Column (Content bullet list)
left_x = Inches(0.8)
add_text_box(slide_3, left_x, Inches(2.0), Inches(5.8), Inches(1.0), 
             "Kami mendesain sistem POS yang berjalan mandiri di browser kasir secara 100% offline, namun terintegrasi cloud secara seamless.", 
             size=15, color=COLOR_WHITE)

bullet_y = Inches(3.2)
# Bullet 1
circle1 = slide_3.shapes.add_shape(MSO_SHAPE.OVAL, left_x, bullet_y, Inches(0.3), Inches(0.3))
circle1.fill.solid()
circle1.fill.fore_color.rgb = COLOR_GREEN
circle1.line.fill.background()
add_text_box(slide_3, left_x + Inches(0.5), bullet_y - Inches(0.05), Inches(5.3), Inches(1.0), 
             "Downtime Tolerant\nOperasional kasir berjalan 100% normal tanpa koneksi internet (baca & tulis database lokal).", size=12, color=COLOR_MUTED)

# Bullet 2
bullet_y2 = bullet_y + Inches(1.1)
circle2 = slide_3.shapes.add_shape(MSO_SHAPE.OVAL, left_x, bullet_y2, Inches(0.3), Inches(0.3))
circle2.fill.solid()
circle2.fill.fore_color.rgb = COLOR_GREEN
circle2.line.fill.background()
add_text_box(slide_3, left_x + Inches(0.5), bullet_y2 - Inches(0.05), Inches(5.3), Inches(1.0), 
             "Sinkronisasi Otomatis & Cerdas\nBegitu koneksi internet terdeteksi kembali, sinkronisasi latar belakang berjalan aman ke database Supabase tanpa interupsi.", size=12, color=COLOR_MUTED)

# Bullet 3
bullet_y3 = bullet_y2 + Inches(1.1)
circle3 = slide_3.shapes.add_shape(MSO_SHAPE.OVAL, left_x, bullet_y3, Inches(0.3), Inches(0.3))
circle3.fill.solid()
circle3.fill.fore_color.rgb = COLOR_GREEN
circle3.line.fill.background()
add_text_box(slide_3, left_x + Inches(0.5), bullet_y3 - Inches(0.05), Inches(5.3), Inches(1.0), 
             "Instalasi Instan & Tanpa Alat Khusus\nDapat dipasang sebagai PWA (Progressive Web App) pada tablet, handphone, maupun laptop kasir langsung dari browser.", size=12, color=COLOR_MUTED)

# Right Column: Visual Diagram Cards
right_x = Inches(7.2)
add_card(slide_3, right_x, Inches(2.2), Inches(5.3), Inches(1.2), 
         "💻 POS Kasir (IndexedDB / Dexie.js)", 
         ["Menampung katalog produk, stok bahan, struk antrean, dan data transaksi secara lokal di memori browser."])

# Connector Arrow Shape
arrow_shape = slide_3.shapes.add_shape(MSO_SHAPE.UP_DOWN_ARROW, right_x + Inches(2.2), Inches(3.6), Inches(0.8), Inches(0.6))
arrow_shape.fill.solid()
arrow_shape.fill.fore_color.rgb = COLOR_LIGHT
arrow_shape.line.fill.background()
add_text_box(slide_3, right_x + Inches(3.2), Inches(3.7), Inches(2.0), Inches(0.4), "Background Sync", size=10, bold=True, color=COLOR_LIGHT)

add_card(slide_3, right_x, Inches(4.4), Inches(5.3), Inches(1.2), 
         "☁️ Supabase PostgreSQL Cloud", 
         ["Sinkronisasi data penjualan real-time, backup otomatis harian, verifikasi lisensi perangkat, dan dashboard analitik owner."])

# --------------------------------------------------------------------------
# SLIDE 4: KEY FEATURES
# --------------------------------------------------------------------------
slide_4 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide_4, "Key Features", "Fitur Komprehensif Siap Operasional UMKM F&B")

# Grid 4 cards
add_card(slide_4, Inches(0.8), Inches(2.2), Inches(2.7), Inches(4.2), 
         "🛒 Kasir & Add-on", 
         [
             "Input pesanan instan.",
             "Manajemen add-on / varian rasa produk (seperti sirup tambahan, espresso shot).",
             "Pemberian diskon manual atau otomatis harian."
         ])

add_card(slide_4, Inches(3.8), Inches(2.2), Inches(2.7), Inches(4.2), 
         "👥 Antrean Visual", 
         [
             "Modul antrean dapur terintegrasi.",
             "Pembagian status visual: Menunggu, Diproses, hingga Siap Disajikan kepada konsumen."
         ])

add_card(slide_4, Inches(6.8), Inches(2.2), Inches(2.7), Inches(4.2), 
         "📦 Bahan Baku", 
         [
             "Manajemen stok resep bahan baku.",
             "Pengurangan bahan secara otomatis saat menu terjual (misal: memotong biji kopi & susu per cup)."
         ])

add_card(slide_4, Inches(9.8), Inches(2.2), Inches(2.7), Inches(4.2), 
         "📊 Laporan & Cloud", 
         [
             "Laporan penjualan, item terlaris, dan omzet harian.",
             "Lisensi aktivasi cloud untuk keamanan B2B SaaS.",
             "Auto backup lokal dan server."
         ])

# --------------------------------------------------------------------------
# SLIDE 5: ARCHITECTURE & TECH STACK
# --------------------------------------------------------------------------
slide_5 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide_5, "Tech Stack", "Kombinasi Teknologi Modern & Skalabilitas Tinggi")

tech_items = [
    ("React 19 & Vite", "Frontend", 
     ["Render komponen UI kasir yang sangat cepat (latensi rendering < 16ms).", 
      "Pemuatan aplikasi awal instan di bawah 1 detik.", 
      "Ukuran bundle ringan untuk menekan konsumsi data seluler."]),
      
    ("Dexie.js & IndexedDB", "Local DB", 
     ["Database transaksional lokal terindeks di browser kasir.", 
      "Kecepatan tulis data transaksi instan (0ms) tanpa menunggu respon jaringan internet.", 
      "Kapasitas penyimpanan data offline hingga gigabyte."]),
      
    ("Supabase Server", "Cloud Backend", 
     ["Backend-as-a-Service berbasis database relasional PostgreSQL.", 
      "Sinkronisasi data real-time, autentikasi multi-kasir, dan penyimpanan foto menu.", 
      "Sistem otentikasi lisensi terpusat."]),
      
    ("PWA (Progressive Web App)", "Deployment", 
     ["Dapat diinstal langsung dari browser ke home screen perangkat kasir.", 
      "Kompatibel penuh di semua OS (Android, iOS, Windows, macOS).", 
      "Mendukung pembaruan aplikasi otomatis tanpa campur tangan user."])
]

for idx, (title, label, bullets) in enumerate(tech_items):
    col_x = Inches(0.8) + Inches(idx * 3.0)
    
    # Graphic Circle for Tech Icon placeholder
    circle = slide_5.shapes.add_shape(MSO_SHAPE.OVAL, col_x + Inches(0.9), Inches(2.0), Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = COLOR_CARD_BG
    circle.line.color.rgb = COLOR_PRIMARY
    circle.line.width = Pt(1.5)
    
    # Text icon simulation
    add_text_box(slide_5, col_x + Inches(0.9), Inches(2.3), Inches(1.0), Inches(0.4), title[:2].upper(), size=14, bold=True, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)
    
    # Card Content
    add_card(slide_5, col_x, Inches(3.2), Inches(2.8), Inches(3.2), title, bullets)

# --------------------------------------------------------------------------
# SLIDE 6: MARKET OPPORTUNITY (WITH NATIVE CHART)
# --------------------------------------------------------------------------
slide_6 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide_6, "Market Opportunity", "Peluang Pasar & Target Segmentasi UMKM F&B")

# Description box left
add_text_box(slide_6, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.8), 
             "Bestari POS membidik sektor ritel makanan dan minuman (F&B) kelas menengah ke bawah yang tersebar di kota tier 2 dan 3 di Indonesia.", 
             size=14, color=COLOR_WHITE)

add_card(slide_6, Inches(0.8), Inches(2.9), Inches(5.5), Inches(1.2), 
         "Total Addressable Market (TAM)", 
         ["Terdapat 64.2 Juta+ unit UMKM di Indonesia, yang merupakan fondasi utama perekonomian nasional."])

add_card(slide_6, Inches(0.8), Inches(4.3), Inches(5.5), Inches(1.2), 
         "Serviceable Addressable Market (SAM)", 
         ["2.4 Juta+ gerai ritel kuliner & F&B aktif (kedai kopi mandiri, warung makan lokal, franchise kecil)."])

add_card(slide_6, Inches(0.8), Inches(5.7), Inches(5.5), Inches(1.2), 
         "Serviceable Obtainable Market (SOM)", 
         ["Target awal mengakuisisi 10.000 merchant aktif dalam jangka waktu 3 tahun pasca-pendanaan."])

# Native Chart Right: Segmentasi Target Merchant (Doughnut Chart)
chart_data = CategoryChartData()
chart_data.categories = ['Kedai Kopi & Kafe Kecil (55%)', 'Kedai Fast-Food Lokal (30%)', 'Franchise Multi-Outlet (15%)']
chart_data.add_series('Proporsi Segmentasi', (0.55, 0.30, 0.15))

chart_x, chart_y, chart_cx, chart_cy = Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.8)
chart = slide_6.shapes.add_chart(
    XL_CHART_TYPE.DOUGHNUT, chart_x, chart_y, chart_cx, chart_cy, chart_data
).chart

chart.has_legend = True
chart.legend.position = XL_LEGEND_POSITION.BOTTOM
chart.legend.font.color.rgb = COLOR_WHITE
chart.legend.font.size = Pt(9)
chart.legend.font.name = FONT_BODY

# Chart Title
add_text_box(slide_6, chart_x, chart_y + Inches(0.1), chart_cx, Inches(0.4), 
             "Fokus Segmentasi Target Merchant Awal", size=13, bold=True, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)

# --------------------------------------------------------------------------
# SLIDE 7: BUSINESS MODEL
# --------------------------------------------------------------------------
slide_7 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide_7, "Business Model", "Model Pendapatan SaaS B2B Skalabel")

# Grid of 3 pricing plans
add_card(slide_7, Inches(0.8), Inches(2.2), Inches(3.6), Inches(4.4), 
         "Starter (Lokal Mandiri)", 
         [
             "Rp 149.000 / Bulan",
             "------------------------",
             "• Maksimal 1 perangkat aktif.",
             "• Akses fitur kasir & manajemen varian penuh.",
             "• Penyimpanan database lokal di browser kasir.",
             "• Ekspor laporan penjualan ke Excel secara offline.",
             "• Dukungan offline-first penuh."
         ])

# Middle card represents the recommended plan (using secondary light accent)
add_card(slide_7, Inches(4.8), Inches(2.2), Inches(3.8), Inches(4.6), 
         "Professional (Cloud Linked)", 
         [
             "Rp 299.000 / Bulan",
             "------------------------",
             "• Sinkronisasi multi-device (kasir & dapur).",
             "• Cloud backup otomatis ke Supabase.",
             "• Modul stok & kalkulator resep bahan baku.",
             "• Dashboard grafik laporan owner (real-time).",
             "• Prioritas update fitur & dukungan support."
         ], accent_color=COLOR_LIGHT)

add_card(slide_7, Inches(8.9), Inches(2.2), Inches(3.6), Inches(4.4), 
         "White-Label (Enterprise)", 
         [
             "Skema Lisensi / Kustom",
             "------------------------",
             "• Domain kustom & logo usaha mandiri.",
             "• Dedicated cloud server database Supabase.",
             "• Integrasi API ke sistem pemesanan luar.",
             "• Akses source code penuh.",
             "• SLA Dukungan 24/7 Premium."
         ])

# --------------------------------------------------------------------------
# SLIDE 8: FINANCIAL PROJECTIONS (WITH NATIVE CHART)
# --------------------------------------------------------------------------
slide_8 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide_8, "Financial Projections", "Proyeksi Pertumbuhan Merchant & Pendapatan ARR")

# Content Left
add_text_box(slide_8, Inches(0.8), Inches(2.0), Inches(5.5), Inches(0.6), 
             "Proyeksi target pencapaian bisnis dalam jangka waktu 3 tahun ke depan:", 
             size=14, color=COLOR_WHITE)

metrics = [
    ("Tahun 1:", "1.500 Merchant Aktif | ARR Rp 3.5 Miliar", "Fokus akuisisi di wilayah perkotaan Jabodetabek."),
    ("Tahun 2:", "5.000 Merchant Aktif | ARR Rp 12.0 Miliar", "Ekspansi pasar ke wilayah Bandung, Surabaya, dan Yogyakarta."),
    ("Tahun 3:", "12.000 Merchant Aktif | ARR Rp 30.0 Miliar", "Kemitraan korporasi F&B multi-outlet di tingkat nasional.")
]

for idx, (yr, val, desc) in enumerate(metrics):
    my = Inches(2.8) + Inches(idx * 1.4)
    add_text_box(slide_8, Inches(0.8), my, Inches(5.5), Inches(0.3), yr, size=13, bold=True, color=COLOR_LIGHT)
    add_text_box(slide_8, Inches(0.8), my + Inches(0.25), Inches(5.5), Inches(0.4), val, size=15, bold=True, color=COLOR_WHITE)
    add_text_box(slide_8, Inches(0.8), my + Inches(0.6), Inches(5.5), Inches(0.4), desc, size=11, color=COLOR_MUTED)

# Native Column Chart Right: ARR Projection
chart_data = CategoryChartData()
chart_data.categories = ['Tahun 1', 'Tahun 2', 'Tahun 3']
chart_data.add_series('ARR (Rp Miliar)', (3.5, 12.0, 30.0))

chart_x, chart_y, chart_cx, chart_cy = Inches(7.0), Inches(2.0), Inches(5.5), Inches(4.5)
chart = slide_8.shapes.add_chart(
    XL_CHART_TYPE.COLUMN_CLUSTERED, chart_x, chart_y, chart_cx, chart_cy, chart_data
).chart

chart.has_legend = False

# Stylize Axis label text colors for dark theme
category_axis = chart.category_axis
category_axis.tick_labels.font.color.rgb = COLOR_WHITE
category_axis.tick_labels.font.size = Pt(10)
category_axis.tick_labels.font.name = FONT_BODY

value_axis = chart.value_axis
value_axis.tick_labels.font.color.rgb = COLOR_WHITE
value_axis.tick_labels.font.size = Pt(10)
value_axis.tick_labels.font.name = FONT_BODY

add_text_box(slide_8, chart_x, chart_y + Inches(0.1), chart_cx, Inches(0.4), 
             "Proyeksi ARR (Annual Recurring Revenue) dalam Rp Miliar", size=11, bold=True, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)

# --------------------------------------------------------------------------
# SLIDE 9: THE ASK / FUNDING ALLOCATION
# --------------------------------------------------------------------------
slide_9 = prs.slides.add_slide(prs.slide_layouts[6])
add_slide_header(slide_9, "The Ask", "Rencana Alokasi Dana Pendanaan Benih")

# Left Column (The Ask)
left_x = Inches(0.8)
add_text_box(slide_9, left_x, Inches(2.2), Inches(5.5), Inches(0.4), "KEBUTUHAN DANA SEED ROUND:", size=12, bold=True, color=COLOR_MUTED)
add_text_box(slide_9, left_x, Inches(2.6), Inches(5.5), Inches(0.9), "Rp 2.000.000.000", size=48, bold=True, color=COLOR_LIGHT)

add_text_box(slide_9, left_x, Inches(3.8), Inches(5.5), Inches(2.5), 
             "Investasi ini direncanakan untuk mempercepat perluasan jangkauan pasar Bestari POS, menyelesaikan modul integrasi pembayaran cashless otomatis (QRIS), serta merekrut tim inti pengembang produk dan pemasaran lapangan.", 
             size=14, color=COLOR_WHITE)

# Right Column (Allocation Progress Bars)
right_x = Inches(7.0)
allocations = [
    ("💻 R&D & Pengembangan Produk (50%)", "Rp 1.0 Miliar", 0.50, 
     "Integrasi pembayaran QRIS dinamis otomatis, e-wallet, dan modul dashboard owner versi mobile iOS & Android."),
     
    ("🚀 Sales & Pemasaran Lapangan (35%)", "Rp 700 Juta", 0.35, 
     "Akuisisi merchant langsung di kota-kota tier 2, kampanye digital terarah, dan promosi komunitas UMKM kuliner."),
     
    ("⚙️ Infrastruktur & Operasional (15%)", "Rp 300 Juta", 0.15, 
     "Biaya server cloud database Supabase, lisensi perangkat lunak keamanan data, operasional legalitas badan usaha.")
]

for idx, (title, amount, ratio, desc) in enumerate(allocations):
    ay = Inches(2.2) + Inches(idx * 1.5)
    add_text_box(slide_9, right_x, ay, Inches(4.0), Inches(0.3), title, size=11, bold=True, color=COLOR_WHITE)
    add_text_box(slide_9, right_x + Inches(4.0), ay, Inches(1.5), Inches(0.3), amount, size=11, bold=True, color=COLOR_LIGHT, align=PP_ALIGN.RIGHT)
    
    # Progress Bar background shape
    bar_bg = slide_9.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, ay + Inches(0.35), Inches(5.5), Inches(0.12))
    bar_bg.fill.solid()
    bar_bg.fill.fore_color.rgb = COLOR_CARD_BG
    bar_bg.line.fill.background()
    
    # Progress Bar fill
    bar_fill = slide_9.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, ay + Inches(0.35), Inches(5.5 * ratio), Inches(0.12))
    bar_fill.fill.solid()
    bar_fill.fill.fore_color.rgb = COLOR_PRIMARY
    bar_fill.line.fill.background()
    
    add_text_box(slide_9, right_x, ay + Inches(0.55), Inches(5.5), Inches(0.5), desc, size=9.5, color=COLOR_MUTED)

# --------------------------------------------------------------------------
# SLIDE 10: PENUTUP & KONTAK
# --------------------------------------------------------------------------
slide_10 = prs.slides.add_slide(prs.slide_layouts[6])
apply_slide_bg(slide_10)

# Large accent container box
container = slide_10.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.5))
container.fill.solid()
container.fill.fore_color.rgb = COLOR_CARD_BG
container.line.color.rgb = COLOR_PRIMARY
container.line.width = Pt(2)

# Closing Text
add_text_box(slide_10, Inches(2.0), Inches(2.2), Inches(9.333), Inches(0.5), 
             "Mari Berkolaborasi Mendigitalisasi UMKM Indonesia!", size=24, bold=True, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)

add_text_box(slide_10, Inches(2.0), Inches(3.0), Inches(9.333), Inches(0.8), 
             "Bergabunglah bersama kami sebagai investor untuk menghadirkan solusi kasir handal bebas downtime bagi jutaan pengusaha kuliner tanah air.", 
             size=14, color=COLOR_WHITE, align=PP_ALIGN.CENTER)

# Contact info dividers
divider = slide_10.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.0), Inches(3.9), Inches(5.333), Inches(0.015))
divider.fill.solid()
divider.fill.fore_color.rgb = COLOR_DARK
divider.line.fill.background()

add_text_box(slide_10, Inches(2.0), Inches(4.2), Inches(9.333), Inches(0.8), 
             "Hubungi Kami untuk Diskusi Lebih Lanjut:\n📧 email: info@bestaripos.com   |   🌐 website: bestaripos.com   |   📍 Jakarta, Indonesia", 
             size=12, color=COLOR_MUTED, align=PP_ALIGN.CENTER)

# Save the generated presentation
output_path = "../Bestari_POS_Pitch_Deck.pptx"
try:
    prs.save(output_path)
    print(f"Presentation saved successfully to: {os.path.abspath(output_path)}")
except Exception as e:
    # If standard parent directory save fails, try saving locally
    output_path = "Bestari_POS_Pitch_Deck.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to local folder: {os.path.abspath(output_path)}")
